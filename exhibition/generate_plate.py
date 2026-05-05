"""Reminq 展示プレート (.pptx) ジェネレータ

生成AIなんでも展示会向けの 1 枚スライド。
13.333" x 7.5" (PPTX 標準ワイドスクリーン = Google Slides 互換)。

使い方:
    python3.12 -m pip install --user python-pptx
    python3.12 exhibition/generate_plate.py

出力:
    exhibition/reminq-ai-exhibition-plate.pptx

Google Slides へのインポート:
    Slides で「ファイル → スライドをインポート → アップロード」で .pptx を選択。
    各テキスト/図形は編集可能なまま取り込まれる。
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu


# ============================================================
# Design tokens — Reminq cool-light
# ============================================================

INK_PRIMARY   = RGBColor(0x0F, 0x17, 0x2A)
INK_SECONDARY = RGBColor(0x47, 0x55, 0x69)
INK_MUTED     = RGBColor(0x94, 0xA3, 0xB8)
BG_BASE       = RGBColor(0xF5, 0xF7, 0xFA)
BG_CARD       = RGBColor(0xFF, 0xFF, 0xFF)
BG_SOFT       = RGBColor(0xEC, 0xF1, 0xF6)
BORDER_SOFT   = RGBColor(0xE2, 0xE8, 0xF0)
ACCENT        = RGBColor(0x25, 0x63, 0xEB)
ACCENT_DEEP   = RGBColor(0x1D, 0x4E, 0xD8)
ACCENT_TINT   = RGBColor(0xDB, 0xE7, 0xFE)
STATUS_TODO   = RGBColor(0x63, 0x66, 0xF1)
STATUS_PROG   = RGBColor(0x0E, 0xA5, 0xE9)
STATUS_DONE   = RGBColor(0x10, 0xB9, 0x81)

FONT_HEAD = "Inter Tight"
FONT_BODY = "Inter"
FONT_JP   = "Noto Sans JP"

ROOT = Path(__file__).resolve().parent
REPO = ROOT.parent
QR_PATH = REPO / "images" / "qr.png"
OUT_PATH = ROOT / "reminq-ai-exhibition-plate.pptx"


# ============================================================
# Helpers
# ============================================================

def add_rect(slide, x, y, w, h, fill, line=None, line_w=0.75, shadow=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    if not shadow:
        shape.shadow.inherit = False
    shape.text_frame.margin_left = 0
    shape.text_frame.margin_right = 0
    shape.text_frame.margin_top = 0
    shape.text_frame.margin_bottom = 0
    return shape


def add_round_rect(slide, x, y, w, h, fill, line=None, line_w=0.75, corner=0.06):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = corner
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    shape.shadow.inherit = False
    return shape


def add_text(
    slide, x, y, w, h, text,
    *, font=FONT_JP, size=12, bold=False, color=INK_PRIMARY,
    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.25,
):
    """Single-paragraph text box."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_runs(
    slide, x, y, w, h, runs,
    *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.25,
):
    """One paragraph, multiple styled runs.

    runs: list of dict(text, font, size, bold, color).
    """
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for spec in runs:
        run = p.add_run()
        run.text = spec["text"]
        run.font.name = spec.get("font", FONT_JP)
        run.font.size = Pt(spec.get("size", 12))
        run.font.bold = spec.get("bold", False)
        run.font.color.rgb = spec.get("color", INK_PRIMARY)
    return box


def add_paragraphs(
    slide, x, y, w, h, paragraphs,
    *, anchor=MSO_ANCHOR.TOP,
):
    """Multi-paragraph textbox.

    paragraphs: list of dict(text, font, size, bold, color, align, line_spacing, space_before).
    """
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    for i, spec in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = spec.get("align", PP_ALIGN.LEFT)
        p.line_spacing = spec.get("line_spacing", 1.3)
        if spec.get("space_before") is not None:
            p.space_before = Pt(spec["space_before"])
        run = p.add_run()
        run.text = spec["text"]
        run.font.name = spec.get("font", FONT_JP)
        run.font.size = Pt(spec.get("size", 12))
        run.font.bold = spec.get("bold", False)
        run.font.color.rgb = spec.get("color", INK_PRIMARY)
    return box


# ============================================================
# Slide composition
# ============================================================

def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Background wash
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, BG_BASE)

    # Top accent bar (thin strip)
    add_rect(slide, 0, 0, prs.slide_width, Inches(0.08), ACCENT)

    # ----- Header row -----
    header_y = Inches(0.42)
    # Brand mark: square + wordmark
    add_round_rect(slide, Inches(0.55), header_y - Inches(0.02),
                   Inches(0.36), Inches(0.36), ACCENT, corner=0.22)
    add_text(slide, Inches(0.62), header_y - Inches(0.04),
             Inches(0.30), Inches(0.40), "R",
             font=FONT_HEAD, size=18, bold=True, color=BG_CARD,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(1.00), header_y - Inches(0.02),
             Inches(2.0), Inches(0.4), "Reminq",
             font=FONT_HEAD, size=22, bold=True, color=INK_PRIMARY)

    # Right side: event eyebrow
    add_round_rect(slide, Inches(9.55), header_y,
                   Inches(3.30), Inches(0.34), BG_CARD,
                   line=BORDER_SOFT, line_w=0.75, corner=0.5)
    add_runs(
        slide, Inches(9.55), header_y, Inches(3.30), Inches(0.34),
        [
            {"text": "●  ", "font": FONT_HEAD, "size": 10, "bold": True, "color": ACCENT},
            {"text": "EXHIBIT  ", "font": FONT_HEAD, "size": 9, "bold": True, "color": INK_MUTED},
            {"text": "生成AIなんでも展示会", "font": FONT_JP, "size": 10, "bold": True, "color": INK_PRIMARY},
        ],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0,
    )

    # ----- Hero area -----
    hero_x = Inches(0.55)
    hero_y = Inches(1.05)

    # Eyebrow
    add_runs(
        slide, hero_x, hero_y, Inches(8.0), Inches(0.30),
        [
            {"text": "AI × ", "font": FONT_HEAD, "size": 11, "bold": True, "color": ACCENT},
            {"text": "Apple リマインダー", "font": FONT_JP, "size": 11, "bold": True, "color": INK_SECONDARY},
        ],
        line_spacing=1.0,
    )

    # Title — two lines via paragraphs
    add_paragraphs(
        slide, hero_x, hero_y + Inches(0.32), Inches(9.5), Inches(1.7),
        [
            {"text": "リマインダーを、", "font": FONT_JP, "size": 40, "bold": True,
             "color": INK_PRIMARY, "line_spacing": 1.15},
            {"text": "AI が読み解く。", "font": FONT_JP, "size": 40, "bold": True,
             "color": ACCENT, "line_spacing": 1.15},
        ],
    )

    # Sub-lead
    add_text(
        slide, hero_x, hero_y + Inches(2.05), Inches(8.6), Inches(0.7),
        "Apple のリマインダーを AI が分析。曖昧で大きなタスクを「動ける単位」に分解し、優先度まで提案する iOS / iPadOS / macOS アプリ。",
        font=FONT_JP, size=13, color=INK_SECONDARY, line_spacing=1.5,
    )

    # ----- Right column: AI demo card (mini) -----
    demo_x = Inches(9.20)
    demo_y = Inches(1.35)
    demo_w = Inches(3.65)
    demo_h = Inches(2.55)

    add_round_rect(slide, demo_x, demo_y, demo_w, demo_h, BG_CARD,
                   line=BORDER_SOFT, corner=0.06)
    # mock window header
    add_rect(slide, demo_x, demo_y, demo_w, Inches(0.30), BG_SOFT)
    add_text(slide, demo_x + Inches(0.18), demo_y + Inches(0.04),
             demo_w - Inches(0.4), Inches(0.22),
             "reminq · ai-breakdown",
             font=FONT_BODY, size=8.5, color=INK_MUTED, line_spacing=1.0)
    # traffic dots
    for i, c in enumerate([RGBColor(0xEF,0x44,0x44), RGBColor(0xF5,0x9E,0x0B), RGBColor(0x10,0xB9,0x81)]):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
            demo_x + Inches(0.18 + i*0.18), demo_y + Inches(0.10),
            Inches(0.10), Inches(0.10))
        dot.fill.solid(); dot.fill.fore_color.rgb = c
        dot.line.fill.background(); dot.shadow.inherit = False

    # prompt
    prompt_y = demo_y + Inches(0.45)
    add_round_rect(slide, demo_x + Inches(0.18), prompt_y,
                   demo_w - Inches(0.36), Inches(0.40), BG_SOFT, corner=0.3)
    add_runs(
        slide, demo_x + Inches(0.30), prompt_y, demo_w - Inches(0.4), Inches(0.40),
        [
            {"text": "▍", "font": FONT_HEAD, "size": 12, "bold": True, "color": ACCENT},
            {"text": " 資料を作る", "font": FONT_JP, "size": 12, "bold": True, "color": INK_PRIMARY},
        ],
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0,
    )

    # arrow / label
    add_runs(
        slide, demo_x + Inches(0.18), prompt_y + Inches(0.50),
        demo_w - Inches(0.36), Inches(0.24),
        [
            {"text": "↓ ", "font": FONT_HEAD, "size": 9, "bold": True, "color": ACCENT},
            {"text": "AI が分解した実行ステップ", "font": FONT_JP, "size": 9, "bold": True, "color": INK_MUTED},
        ],
        line_spacing=1.0,
    )

    # generated steps with priority chips
    steps = [
        ("目次を 5 行で書き出す", "High", STATUS_TODO),
        ("参考資料を 3 件集める", "Mid",  STATUS_PROG),
        ("結論スライドを 1 枚先に作る", "High", STATUS_TODO),
        ("空白スライドにタイトルだけ並べる", "Low", STATUS_DONE),
    ]
    step_y = prompt_y + Inches(0.85)
    for i, (txt, prio, color) in enumerate(steps):
        y = step_y + Inches(0.32 * i)
        # checkbox
        add_round_rect(slide, demo_x + Inches(0.20), y + Inches(0.05),
                       Inches(0.16), Inches(0.16), BG_CARD,
                       line=BORDER_SOFT, line_w=1.0, corner=0.2)
        # text
        add_text(slide, demo_x + Inches(0.42), y, Inches(2.40), Inches(0.28),
                 txt, font=FONT_JP, size=9, color=INK_PRIMARY,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
        # priority chip
        chip_w = Inches(0.55)
        chip_x = demo_x + demo_w - Inches(0.22) - chip_w
        add_round_rect(slide, chip_x, y + Inches(0.04),
                       chip_w, Inches(0.20), color, corner=0.5)
        add_text(slide, chip_x, y + Inches(0.04),
                 chip_w, Inches(0.20), prio,
                 font=FONT_HEAD, size=8, bold=True, color=BG_CARD,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

    # ----- Three feature cards -----
    cards_y = Inches(4.05)
    card_w = Inches(4.05)
    card_h = Inches(2.55)
    gap = Inches(0.13)
    margin_l = Inches(0.55)

    cards = [
        {
            "num": "01",
            "title": "AI タスク分解",
            "what": "「資料を作る」のような曖昧で大きな1行を、実行可能な小ステップ群に展開。",
            "where": "入力タスク文字列 → OpenAI Chat Completions API → 構造化レスポンス",
            "effect": "「最初の一歩」が常に提示される — 着手率が上がり、つまずく時点を後ろにずらせる。",
        },
        {
            "num": "02",
            "title": "優先度ラベリング",
            "what": "分解した各ステップに High / Mid / Low の優先度を AI が同時付与。",
            "where": "分解レスポンスの priority フィールドを 1 回の生成で同時出力",
            "effect": "フラットなリストでは「自分で考えて」だった粒度・順序判断に AI が下案を提示。",
        },
        {
            "num": "03",
            "title": "Human-in-the-loop 設計",
            "what": "提案はチェックボックスで採用したものだけが Apple リマインダーへ静かに反映。",
            "where": "AI = 提案担当 / 人間 = 最終決定。BYO API キーで端末→OpenAI 直通、中継なし。",
            "effect": "AI 出力をそのまま受けない設計。誤りに強く、データは Apple リマインダーに留まる。",
        },
    ]

    for i, c in enumerate(cards):
        x = margin_l + (card_w + gap) * i

        # card body
        add_round_rect(slide, x, cards_y, card_w, card_h, BG_CARD,
                       line=BORDER_SOFT, corner=0.04)
        # left accent bar
        add_rect(slide, x, cards_y, Inches(0.06), card_h, ACCENT)

        pad = Inches(0.28)
        # number + title row
        add_runs(
            slide, x + pad, cards_y + Inches(0.22),
            card_w - pad*2, Inches(0.34),
            [
                {"text": c["num"], "font": FONT_HEAD, "size": 11, "bold": True, "color": ACCENT},
                {"text": "   " + c["title"], "font": FONT_JP, "size": 16, "bold": True, "color": INK_PRIMARY},
            ],
            line_spacing=1.0, anchor=MSO_ANCHOR.MIDDLE,
        )

        # divider
        add_rect(slide, x + pad, cards_y + Inches(0.66),
                 card_w - pad*2, Emu(6350), BORDER_SOFT)

        # what
        add_text(
            slide, x + pad, cards_y + Inches(0.78),
            card_w - pad*2, Inches(0.7),
            c["what"],
            font=FONT_JP, size=10.5, color=INK_PRIMARY, line_spacing=1.45,
        )

        # where (label + body)
        where_y = cards_y + Inches(1.50)
        add_text(
            slide, x + pad, where_y,
            card_w - pad*2, Inches(0.22),
            "WHERE  AI を使う場所",
            font=FONT_HEAD, size=8, bold=True, color=ACCENT, line_spacing=1.0,
        )
        add_text(
            slide, x + pad, where_y + Inches(0.22),
            card_w - pad*2, Inches(0.45),
            c["where"],
            font=FONT_JP, size=9.5, color=INK_SECONDARY, line_spacing=1.4,
        )

        # effect (label + body)
        eff_y = cards_y + Inches(2.05)
        add_text(
            slide, x + pad, eff_y,
            card_w - pad*2, Inches(0.22),
            "EFFECT  効果",
            font=FONT_HEAD, size=8, bold=True, color=ACCENT, line_spacing=1.0,
        )
        add_text(
            slide, x + pad, eff_y + Inches(0.22),
            card_w - pad*2, Inches(0.45),
            c["effect"],
            font=FONT_JP, size=9.5, color=INK_SECONDARY, line_spacing=1.4,
        )

    # ----- Footer -----
    foot_y = Inches(6.85)

    # left: QR + URL
    qr_size = Inches(0.55)
    if QR_PATH.exists():
        slide.shapes.add_picture(str(QR_PATH), Inches(0.55), foot_y - Inches(0.03),
                                 width=qr_size, height=qr_size)
    add_paragraphs(
        slide, Inches(1.22), foot_y - Inches(0.04), Inches(4.6), Inches(0.6),
        [
            {"text": "詳しくは サポートサイトへ", "font": FONT_JP, "size": 8.5, "bold": True,
             "color": INK_MUTED, "line_spacing": 1.1},
            {"text": "tichise.github.io/reminq-support", "font": FONT_BODY, "size": 11, "bold": True,
             "color": INK_PRIMARY, "line_spacing": 1.1},
        ],
    )

    # center: pill row
    pills = [
        ("iOS 26+", INK_PRIMARY),
        ("iPadOS 26+", INK_PRIMARY),
        ("macOS 26+", INK_PRIMARY),
        ("完全無料", ACCENT_DEEP),
    ]
    pill_x = Inches(6.10)
    pill_y = foot_y + Inches(0.10)
    for label, color in pills:
        w = Inches(1.15) if "26+" in label else Inches(0.95)
        add_round_rect(slide, pill_x, pill_y, w, Inches(0.32),
                       BG_CARD, line=BORDER_SOFT, corner=0.5)
        add_text(slide, pill_x, pill_y, w, Inches(0.32),
                 label,
                 font=FONT_HEAD if "26+" in label else FONT_JP,
                 size=9, bold=True, color=color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        pill_x = pill_x + w + Inches(0.08)

    # right-bottom: tagline
    add_text(
        slide, Inches(10.70), foot_y + Inches(0.10),
        Inches(2.3), Inches(0.32),
        "Reminders, read by AI.",
        font=FONT_HEAD, size=10, bold=True, color=INK_MUTED,
        align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0,
    )

    return prs


def main() -> None:
    prs = build_presentation()
    prs.save(str(OUT_PATH))
    print(f"wrote: {OUT_PATH.relative_to(REPO)}")


if __name__ == "__main__":
    main()
